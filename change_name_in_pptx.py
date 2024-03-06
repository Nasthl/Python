from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_PATTERN_TYPE
from pptx.enum.dml import MSO_FILL
import os
import glob
import datetime
import shutil
#Makro do zmiany nazwy autora w wszystkich pptx w folderze wraz z poprawą rewizji na wszystkich slajdach i slajdzie głównym

# zmienia fragment stringa
# zwraca ilosc zmian
def zmien_fragment_stringa_w_pptx(old_pptx,new_pptx,old_str,new_str):
    ppt = Presentation(old_pptx)
    iterator=0
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if (run.text.find(old_str)) != -1:
                            run.text = run.text.replace(old_str, new_str)
                            iterator+=1
    if iterator:
        ppt.save(new_pptx)
    return iterator

# zmienia caly string
# zwraca ilosc zmian
def zmien_rewizje_w_slajdach(old_pptx,new_pptx,old_str,new_str):
    ppt = Presentation(old_pptx)
    iterator=0
    liczba_slajdow=0
    for slide in ppt.slides:
        liczba_slajdow+=1
        for shape in slide.shapes:
            if shape.has_text_frame:

                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text=''
                    for run in paragraph.runs:
                        if (run.text.find(old_str)) != -1:
                            paragraph_text+=run.text
                    #sprawdz czy dlugosc paragrafu rowna sie dlugosci ciagu i jest pattern jako tlo
                    #tu ciag ma jeden znak wiec podmieniam ostatnia i jedyna iteracje run
                    if len(paragraph_text)==len(new_str) and shape.fill.type==2 and 5900000<shape.top<6200000 and 2500000<shape.left<2800000:
                        #print(shape.left,shape.top)
                        run.text = run.text.replace(old_str, new_str)
                        iterator += 1

    print(f'W {new_pptx} podmieniono {old_str} na {new_str} {iterator} razy')
    if iterator!=liczba_slajdow:
        print(f'Liczba zmian NIE ROWNA sie liczbie slajdow: zmiany {iterator}, slajdy {liczba_slajdow}')
    ppt.save(new_pptx)

def dopisz_rewizje_to_tabeli(old_pptx,new_pptx,rev,comment):
    ppt = Presentation(old_pptx)
    slajd=ppt.slides[0]
    for shape in slajd.shapes:
        if shape.has_table:
            tabela=shape.table
            for i in range(len(tabela.rows)):
                if tabela.cell(i, 1).text == "":
                    tabela.cell(i, 1).text= str(rev)
                    tabela.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(10)
                    tabela.cell(i, 1).text_frame.paragraphs[0].font.name = "VW Head Office"
                    for paragraph in tabela.cell(i, 1).text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                    tabela.cell(i, 2).text = comment
                    tabela.cell(i, 2).text_frame.paragraphs[0].font.size = Pt(10)
                    tabela.cell(i, 2).text_frame.paragraphs[0].font.name = "Arial (Body)"
                    tabela.cell(i, 3).text = str(datetime.datetime.today().strftime('%d.%m.%Y'))
                    tabela.cell(i, 3).text_frame.paragraphs[0].font.size = Pt(10)
                    tabela.cell(i, 3).text_frame.paragraphs[0].font.name = "VW Head Office"
                    tabela.cell(i, 4).text = "CHS"
                    tabela.cell(i, 4).text_frame.paragraphs[0].font.size = Pt(10)
                    tabela.cell(i, 4).text_frame.paragraphs[0].font.name = "VW Head Office"
                    break

    ppt.save(new_pptx)

def zmien_rewizje(old_pptx_name,temp_pptx,old_path,new_path, comment):
    # zmien nazwe ppt
    rewizja = int(old_pptx_name[old_pptx_name.find('_A0') + 3])
    new_pptx_name = old_pptx_name[:old_pptx_name.find('_A0') + 3] + str(rewizja + 1) + old_pptx_name[old_pptx_name.find('_A0') + 4:]
    new_pptx_name=new_pptx_name.replace(old_path,new_path)
    #print(new_pptx_name)
    # zmien rewizje na kazdym slajdzie
    zmien_rewizje_w_slajdach(temp_pptx, new_pptx_name, str(rewizja), str(rewizja + 1))
    # dopisz komentarz do tabeli
    dopisz_rewizje_to_tabeli(new_pptx_name, new_pptx_name, rewizja + 1, comment)

def utworz_strukture_folderowa(old_path, new_path):
    for sciezka, podfoldery, pliki in os.walk(old_path):
        nowa_sciezka=os.path.join(new_path, os.path.relpath(sciezka, old_path))
        if not os.path.exists(nowa_sciezka):
            os.makedirs(nowa_sciezka)

path='C:\\Users\\tomasz.korzyniec\\Downloads\\Plany\\'
new_path='C:\\Users\\tomasz.korzyniec\\Downloads\\Plany_wynik\\'
old_name="O. Voigt"
old_name2="O.Voigt"
new_name="M. Štulajter"
ppt_temp=f'{path}\\temp.pptx'
komentarz='O. Voigt wechselt zu M. Štulajter / O. Voigt sa zmenil na M. Štulajter'

utworz_strukture_folderowa(path,new_path)
for root, dirs, files in os.walk(path):
    pptx_files=glob.glob(os.path.join(root,"*.pptx"))
    for ppt_file in pptx_files:
        ppt = Presentation(ppt_file)
        ppt.save(ppt_temp)
        #Zbior zmian
        zmiana=0
        zmiana+=zmien_fragment_stringa_w_pptx(ppt_temp ,ppt_temp,old_name,new_name)
        zmiana+=zmien_fragment_stringa_w_pptx(ppt_temp, ppt_temp, old_name2, new_name)
        #Jesli zastapila zmiana, zmien rewizje
        if zmiana:
            zmien_rewizje(ppt_file,ppt_temp,path, new_path, komentarz)
        os.remove(ppt_temp)

