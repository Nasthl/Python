
import os
import glob
from pptx import Presentation
import datetime
import shutil

#Makro do znajdowania najnowszych wersji dokumentacji i łączenia w jedno

def copy_shapes(slide, new_slide):
    shape_to_transfer=slide.shapes
    for shape in shape_to_transfer:
        new_shape=new_slide.shapes.add_shape(shape,left=shape.left,
                                                    top=shape.top,
                                                    width=shape.width,
                                                    height=shape.height)




path='C:\\Users\\tkdar\\Downloads\\AnlagendokuMergeTest\\'
path_master='C:\\Users\\tkdar\\Downloads\\AnlagendokuMergeTest\\wzor.pptx'
#new_path='C:\\Users\\tomasz.korzyniec\\Downloads\\Plany_wynik\\'
print('Test')
plans=['A33_Punktplan','A33_RESplan','A37_HSNplan', 'A37_Clinchplan','A39_Klebeplan','A38_Bolzenplan','A40_Schraubenplan']
#lista wszystkich planow
all_pptx_files=[]
for root, dirs, files in os.walk(path):
    pptx_files=glob.glob(os.path.join(root,"*plan*.pptx"))
    for pptx_file in pptx_files:
        all_pptx_files.append(pptx_file)

#usuniecie starszych wersji
temp_all_pptx_files=all_pptx_files[:]
for all_pptx_file in all_pptx_files:
    if int(all_pptx_file[all_pptx_file.find('_A0')+3])>0:
        temp_rev=(int(all_pptx_file[all_pptx_file.find('_A0')+3]))
        for temp in range(temp_rev,0,-1):
            temp_name=all_pptx_file[all_pptx_file.find('_A0')-9:all_pptx_file.find('_A0')+3]+str(temp-1)+all_pptx_file[all_pptx_file.find('_A0')+4:all_pptx_file.find('_A0')+8]
            for temp_all_pptx_file in reversed(temp_all_pptx_files):
                if temp_name in temp_all_pptx_file:
                    temp_all_pptx_files.remove(temp_all_pptx_file)

all_pptx_files=[]
all_pptx_files=temp_all_pptx_files[:]
master_ppt=Presentation(path_master)
for plan in plans:
    master_ppt.save(path+plan+'_A00.pptx')
#podzial na poszczegolne plany (po wyczyszczeniu starych rewizji mozna odrazu przypisywac)
for all_pptx_file in all_pptx_files:
    check_plan = 0
    for plan in plans:
        if all_pptx_file.find(plan)>0:
            check_plan=1
            source_presentation=Presentation(all_pptx_file)
            target_presentation=Presentation(path+plan+'_A00.pptx')
            iterator=0
            for slide in source_presentation.slides:
                iterator+=1
                if iterator>1:
                    new_slide=target_presentation.slides.add_slide(slide.slide_layout)
                    copy_shapes(slide, new_slide)

                #for shape in slide.shapes:
                    #new_shape=new_slide.shapes.add_copy(shape)
            target_presentation.save(path+plan+'_A00.pptx')
            #
    if not(check_plan):
        print('Nie znaleziono ', all_pptx_file)

